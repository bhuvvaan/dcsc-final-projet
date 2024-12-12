import cv2
import os
import numpy as np
from skimage.metrics import structural_similarity as ssim
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import base64
import io
import logging
import os
import redis
from minio import Minio
import glob
import uuid
import json

import psycopg2
from psycopg2.extensions import ISOLATION_LEVEL_AUTOCOMMIT



minioHost = os.getenv("MINIO_HOST") or "localhost:9000"
minioUser = os.getenv("MINIO_USER") or "rootuser"
minioPasswd = os.getenv("MINIO_PASSWD") or "rootpass123"
print(f"Getting minio connection now for host {minioHost}!")

MINIO_CLIENT = None
BUCKET_NAME = "user-uploaded-video-bucket"
PPT_BUCKET_NAME = "ppt-result"
FILE_NAME = 'new_name.mp4'

try:
    MINIO_CLIENT = Minio(minioHost, access_key=minioUser, secret_key=minioPasswd, secure=False)
except Exception as exp:
    print(f"Exception raised in worker loop: {str(exp)}")


# Connect to the default database
conn = psycopg2.connect(
    host="localhost",   # or the IP of your service if not using port-forwarding
    port=5432,
    user="postgres",
    password="my-password",  # Replace with your password
    database="postgres"  # Default database in PostgreSQL
)


# Set the isolation level to AUTOCOMMIT
conn.set_isolation_level(ISOLATION_LEVEL_AUTOCOMMIT)

# Create a cursor object
cur = conn.cursor()



def create_ppt_from_images(image_folder, output_ppt,id):

    prs = Presentation()

    image_files = sorted(os.listdir(image_folder))

    # Loop through all files in the sorted list
    for image_file in image_files:
        # Check if the file is an image (you can add more extensions if needed)
        if image_file.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            image_path = os.path.join(image_folder, image_file)

            # Add a slide with blank layout
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Add the image to the slide
            img = slide.shapes.add_picture(image_path, Inches(0), Inches(0), height=Inches(5))

    # Save the PowerPoint presentation
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)


    # Upload the presentation to MinIO
    ppt_filename = str(id) + "_presentation.pptx"  # The name of the file to be saved in MinIO

    try:
        # Upload the file to MinIO
        MINIO_CLIENT.put_object(PPT_BUCKET_NAME, ppt_filename, ppt_stream, len(ppt_stream.getvalue()))
        print(f"PPT successfully uploaded to MinIO as {ppt_filename}")
    except Exception as e:
        print(f"Error occurred while uploading to MinIO: {e}")



def capture_frames(output_dir, interval,id):
    
    # Step 1: Download video from MinIO to memory (BytesIO)
    video_data = MINIO_CLIENT.get_object(BUCKET_NAME, FILE_NAME)
    video_bytes = video_data.read()  # Read the video data

    # Create a BytesIO object to hold the video in memory
    video_stream = BytesIO(video_bytes)

    # Step 2: Convert BytesIO stream into a format that OpenCV can read
    # OpenCV can read videos from a stream using the VideoCapture class.
    video_stream.seek(0)  # Ensure the stream is at the beginning
    nparr = np.frombuffer(video_stream.read(), np.uint8)
    cap = cv2.VideoCapture()
    cap.open(nparr)

    if not cap.isOpened():
        print(f"Error: Unable to open video file'.")
        return

    # Get video properties
    fps = int(cap.get(cv2.CAP_PROP_FPS))  # Frames per second
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))  # Total number of frames
    duration = total_frames / fps  # Duration of the video in seconds

    print(f"Video Details: FPS={fps}, Total Frames={total_frames}, Duration={duration:.2f}s")

    # Calculate the frame interval based on the given time interval
    frame_interval = fps * interval

    frame_count = 0
    saved_frame_count = 0
    previous_file = None

    while True:
        ret, frame = cap.read()
        if not ret:
            break

        # Save frames at the specified interval
        if frame_count % frame_interval == 0:
            frame_time = frame_count / fps
            output_file = os.path.join(output_dir, f"frame_{saved_frame_count:04d}_{frame_time:.2f}s.jpg")
            
            next_file = os.path.join(output_dir, f"frame_{saved_frame_count + 1:04d}_{frame_time + interval:.2f}s.jpg")
            # print(previous_file)
            # print
            if previous_file and saved_frame_count >= 1:
                print("entered if")
                frame1_path = previous_file

                frame1 = cv2.imread(frame1_path)
                frame2 = frame

                if frame1 is None or frame2 is None:
                    print("Error: One of the frames could not be loaded.")
                else:
                    print("Entered else")
                    mse_val, ssim_val, hist_corr = compare_frames(frame1, frame2)

                
                    if hist_corr <= 0.998 or ssim_val <= 0.97:
                    # Save the current frame
                        cv2.imwrite(output_file, frame)
                        print(f"Saved: {output_file}")
                        saved_frame_count += 1
                        previous_file = output_file
                    
                    # Update previous file for the next iteration
            if saved_frame_count < 1:
                cv2.imwrite(output_file, frame)
                saved_frame_count += 1
                previous_file = output_file

        frame_count += 1

    # Release the video capture object
    cap.release()
    print("Frame capture complete.")

def mse(imageA, imageB):
    """Calculate the Mean Squared Error (MSE) between two images."""
    # The 'Mean Squared Error' between the two images is the sum of the squared differences between the pixel intensities.
    err = np.sum((imageA.astype("float") - imageB.astype("float")) ** 2)
    err /= float(imageA.shape[0] * imageA.shape[1])
    return err


def compare_histograms(imageA, imageB):
    """Compare histograms of two images."""
    histA = cv2.calcHist([imageA], [0], None, [256], [0, 256])
    histB = cv2.calcHist([imageB], [0], None, [256], [0, 256])
    # Normalize the histograms
    histA = cv2.normalize(histA, histA).flatten()
    histB = cv2.normalize(histB, histB).flatten()
    # Compute correlation between histograms
    return cv2.compareHist(histA, histB, cv2.HISTCMP_CORREL)

def compare_frames(frameA, frameB):
    """Compare two frames using different similarity metrics."""
    # Resize frames to the same size if they differ
    if frameA.shape != frameB.shape:
        frameB = cv2.resize(frameB, (frameA.shape[1], frameA.shape[0]))

    # Convert to grayscale for MSE and SSIM comparisons
    grayA = cv2.cvtColor(frameA, cv2.COLOR_BGR2GRAY)
    grayB = cv2.cvtColor(frameB, cv2.COLOR_BGR2GRAY)

    # Compute the MSE between the two images
    mse_value = mse(grayA, grayB)
    print(f"Mean Squared Error (MSE): {mse_value}")

    # Compute the SSIM between the two images
    ssim_value, _ = ssim(grayA, grayB, full=True)
    print(f"Structural Similarity Index (SSIM): {ssim_value}")

    # Compute histogram comparison (correlation)
    histogram_correlation = compare_histograms(frameA, frameB)
    print(f"Histogram Correlation: {histogram_correlation}")

    return mse_value, ssim_value, histogram_correlation

def read_video(FILE_NAME):
    video_path = MINIO_CLIENT.get_object(BUCKET_NAME, FILE_NAME)
    output_dir = ""
    interval = 4
    capture_frames(video_path, output_dir, interval)

    image_folder = output_dir
    # add dynamic name
    output_ppt = r"images_presentation.pptx"

    create_ppt_from_images(image_folder, output_ppt)

def update_postgres(id):
    video_name = str(id) + '_video.mp4'
    video_url = MINIO_CLIENT.presigned_get_object(BUCKET_NAME, video_name)
    print(f"Access the file at: {video_url}")

    ppt_name = str(id) + "_presentation.pptx" 
    ppt_url = MINIO_CLIENT.presigned_get_object(PPT_BUCKET_NAME,FILE_NAME)

    # Define the INSERT query
    insert_query = '''
    INSERT INTO video_details (video_id, video_link_minio, ppt_link) 
    VALUES (%s, %s, %s);
    '''

    data = [id,video_url,ppt_url]

    cur.execute(insert_query, data)
    print("Row inserted successfully!")

    # Commit the transaction
    conn.commit()

    # Close the cursor and connection
    cur.close()
    conn.close()


def capture_frames_2(video_path, output_dir, interval,id):
    
    # Check if the video file exists
    if not os.path.exists(video_path):
        print(f"Error: Video file '{video_path}' does not exist.")
        return

    # Create the output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)

    # Load the video file
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        print(f"Error: Unable to open video file '{video_path}'.")
        return

    # Get video properties
    fps = int(cap.get(cv2.CAP_PROP_FPS))  # Frames per second
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))  # Total number of frames
    duration = total_frames / fps  # Duration of the video in seconds

    print(f"Video Details: FPS={fps}, Total Frames={total_frames}, Duration={duration:.2f}s")

    # Calculate the frame interval based on the given time interval
    frame_interval = fps * interval

    frame_count = 0
    saved_frame_count = 0
    previous_file = None

    while True:
        ret, frame = cap.read()
        if not ret:
            break

        # Save frames at the specified interval
        if frame_count % frame_interval == 0:
            frame_time = frame_count / fps
            output_file = os.path.join(output_dir, f"frame_{saved_frame_count:04d}_{frame_time:.2f}s.jpg")
            
            next_file = os.path.join(output_dir, f"frame_{saved_frame_count + 1:04d}_{frame_time + interval:.2f}s.jpg")
            # print(previous_file)
            # print
            if previous_file and saved_frame_count >= 1:
                print("entered if")
                frame1_path = previous_file

                frame1 = cv2.imread(frame1_path)
                frame2 = frame

                if frame1 is None or frame2 is None:
                    print("Error: One of the frames could not be loaded.")
                else:
                    print("Entered else")
                    mse_val, ssim_val, hist_corr = compare_frames(frame1, frame2)

                    if hist_corr <= 0.998 or ssim_val <= 0.97:
                    # Save the current frame
                        cv2.imwrite(output_file, frame)
                        print(f"Saved: {output_file}")
                        saved_frame_count += 1
                        previous_file = output_file
                    
                    # Update previous file for the next iteration
            if saved_frame_count < 1:
                cv2.imwrite(output_file, frame)
                saved_frame_count += 1
                previous_file = output_file

            
            

        frame_count += 1

    # Release the video capture object
    cap.release()
    print("Frame capture complete.")

if __name__ == "__main__":
    video_path = r"D:\Masters\Fall 2024\data_center_scale_computing\Project\dcsc_final_project\videoplayback.mp4"  # Replace with your video file path
    output_dir = r"D:\Masters\Fall 2024\data_center_scale_computing\Project\dcsc_final_project\frames_2"  # Directory to save frames
    interval = 3  # Capture frames every 5 seconds

    id = 999777
    #id needs to come from redis.

    capture_frames_2(video_path, output_dir, interval,id)

    image_folder = output_dir
    output_ppt = r"D:\Masters\Fall 2024\data_center_scale_computing\Project\dcsc_final_project\images_presentation.pptx"

    create_ppt_from_images(image_folder, output_ppt,id)
    update_postgres(id)

# FILE_NAME = 'new_name.mp4'









